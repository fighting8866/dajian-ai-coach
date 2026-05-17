import re

try:
    from pptx import Presentation
    PPT_AVAILABLE = True
except ImportError:
    PPT_AVAILABLE = False


class PPTService:
    """PPT 文本解析服务（第一阶段：按页提取文本）。"""

    def extract_text_by_slide(self, file_path: str) -> dict:
        """提取整份 PPT 的全文与逐页文本。

        Returns:
            {
              "full_text": "...",
              "slides": [{"page": 1, "text": "..."}, ...]
            }
        """
        if not PPT_AVAILABLE:
            raise ImportError("后端缺少 python-pptx 依赖，请在 backend 目录执行：pip install python-pptx")

        prs = Presentation(file_path)
        slides = []
        full_text_parts = []

        for idx, slide in enumerate(prs.slides, start=1):
            page_text = self._extract_slide_text(slide)
            blocks = self._extract_slide_blocks(slide)
            slides.append(
                {
                    "page": idx,
                    "text": page_text,
                    "blocks": blocks,
                }
            )
            if page_text:
                full_text_parts.append(page_text)

        return {
            "full_text": "\n".join(full_text_parts).strip(),
            "slides": slides
        }

    def parse_ppt(self, file_path: str) -> list[dict]:
        """兼容旧接口：返回 page_index/title/keywords，用于现有匹配流程。"""
        parsed = self.extract_text_by_slide(file_path)
        pages = []
        for slide in parsed["slides"]:
            text = slide.get("text", "")
            title = self._pick_title_from_text(text)
            pages.append({
                "page_index": slide["page"],
                "title": title,
                "keywords": self._extract_keywords(text)
            })
        return pages

    def _iter_shape_text_chunks(self, slide) -> list[tuple[str, str]]:
        """按形状顺序产出 (block_type, chunk_text)；与历史一致，同一 shape 可有 text 与 table 两块。"""
        chunks: list[tuple[str, str]] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = (shape.text or "").strip()
                if text:
                    chunks.append(("text", text))
            if hasattr(shape, "has_table") and shape.has_table:
                for row in shape.table.rows:
                    cell_parts: list[str] = []
                    for cell in row.cells:
                        cell_text = (cell.text or "").strip()
                        if cell_text:
                            cell_parts.append(cell_text)
                    if cell_parts:
                        chunks.append(("table", " | ".join(cell_parts)))
        return chunks

    def _extract_slide_blocks(self, slide) -> list[dict]:
        """结构化块：保留幻灯片内形状顺序，供文档理解 / 后续评分。"""
        return [{"type": typ, "text": txt} for typ, txt in self._iter_shape_text_chunks(slide)]

    def _extract_slide_text(self, slide) -> str:
        """提取单页文本（与历史行为一致：块拼接后按内容去重、保序）。"""
        pieces = [chunk[1] for chunk in self._iter_shape_text_chunks(slide)]
        deduped: list[str] = []
        seen: set[str] = set()
        for item in pieces:
            if item in seen:
                continue
            deduped.append(item)
            seen.add(item)
        return "\n".join(deduped).strip()

    def _pick_title_from_text(self, text: str) -> str:
        """从整页文本中猜测标题（优先第一行非空）。"""
        if not text:
            return ""
        first_line = next((line.strip() for line in text.splitlines() if line.strip()), "")
        return first_line[:80]

    def _extract_keywords(self, text: str) -> list:
        """从文本中提取关键词（用于现有匹配逻辑）。"""
        separators = [" ", "\n", "\t", "，", "。", "；", "、", ",", ".", ";", "：", ":"]
        for sep in separators:
            text = text.replace(sep, " ")
        words = text.split()

        filtered_words = []
        stop_words = {
            "的", "了", "是", "在", "有", "和", "就", "不", "人", "都", "一", "一个",
            "上", "也", "很", "到", "说", "要", "去", "你", "会", "着", "没有", "看", "好", "自己", "这"
        }

        for word in words:
            if len(word) < 2:
                continue
            if word in stop_words:
                continue
            if re.match(r"^\d+$", word):
                continue
            filtered_words.append(word)

        unique_words = list(dict.fromkeys(filtered_words))
        return unique_words[:12] if len(unique_words) > 12 else unique_words